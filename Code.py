# Great work !! , here are some important python functions to make a presentation .

from pptx import Presentation  # Main class to work with PowerPoint presentations
from pptx.util import Inches, Pt, Emu  # Utility classes for measurements
from pptx.enum.shapes import MSO_SHAPE  # Enumeration for shape types
from pptx.chart.data import CategoryChartData  # Class for creating chart data
from pptx.enum.chart import XL_LEGEND_POSITION, XL_AXIS_CROSSES, XL_TICK_LABEL_POSITION  # Enums for chart options

# Presentation Creation and Manipulation
Presentation()  # Creates a new presentation object or opens an existing one

# Adding Slides
Presentation.slides.add_slide(layout)  # Adds a new slide to the presentation with a specific layout

# Accessing and Modifying Slides
Presentation.slides  # List of slides in the presentation
Slide.shapes  # List of shapes on a slide

# Adding and Accessing Shapes
Slide.shapes.add_shape(auto_shape_type, left, top, width, height)  # Adds a shape to a slide
Slide.shapes.add_textbox(left, top, width, height)  # Adds a textbox to a slide
Shape.text  # Gets or sets the text of a shape
Shape.text_frame  # Gets the text frame of a shape
TextFrame.paragraphs  # List of paragraphs in a text frame
Paragraph.text  # Gets or sets the text of a paragraph

# Adding and Accessing Pictures
Slide.shapes.add_picture(image_path, left, top, width, height)  # Adds a picture to a slide

# Adding and Accessing Tables
Slide.shapes.add_table(rows, cols, left, top, width, height)  # Adds a table to a slide
Table.cell(row, col)  # Accesses a specific cell in a table

# Formatting
Shape.fill.solid()  # Sets the fill of a shape to a solid color
Shape.fill.pattern()  # Sets the fill of a shape to a pattern
Shape.line.color.rgb  # Sets the color of the line of a shape
Paragraph.font.size  # Sets the font size of a paragraph
Paragraph.font.bold  # Sets the font to bold
Shape.left  # Gets or sets the left position of a shape
Shape.top  # Gets or sets the top position of a shape
Shape.width  # Gets or sets the width of a shape
Shape.height  # Gets or sets the height of a shape
Shape.rotation  # Gets or sets the rotation of a shape

# Utility Classes
Inches(value)  # Converts inches to EMUs (English Metric Units)
Pt(value)  # Converts points to EMUs
Emu(value)  # Converts EMUs to EMUs (useful for exact measurements)

# Working with Slide Layouts
Presentation.slide_layouts  # List of available slide layouts
SlideLayout.shapes  # Shapes available in a slide layout

# Working with Text Frames
TextFrame.auto_size  # Sets the text frame to auto-size
TextFrame.margin_left  # Sets the left margin of the text frame
TextFrame.margin_right  # Sets the right margin of the text frame
TextFrame.margin_top  # Sets the top margin of the text frame
TextFrame.margin_bottom  # Sets the bottom margin of the text frame

# Working with Charts
Chart = Slide.shapes.add_chart(chart_type, x, y, cx, cy, data)  # Adds a chart to a slide

# Working with Hyperlinks
Shape.click_action.hyperlink.address  # Sets or gets the hyperlink address for a shape

# Working with Slide Master and Layouts
Presentation.slide_master  # Accesses the master slide for a presentation
SlideMaster.placeholders  # List of placeholders in the master slide
SlideMaster.placeholders[index]  # Accesses a specific placeholder in the master slide

# Working with Slide Notes
Slide.notes_slide  # Accesses the notes slide associated with the slide
NotesSlide.notes_text_frame  # Gets or sets the text frame of the notes slide

# Adding and Formatting Shapes
Shape.text_frame.text  # Gets or sets the entire text in a shape
Shape.text_frame.paragraphs[index]  # Accesses a specific paragraph in the text frame
Paragraph.font.name  # Gets or sets the font name of a paragraph
Paragraph.font.color.rgb  # Gets or sets the font color of a paragraph

# Working with Lists
TextFrame.text_frame.add_paragraph()  # Adds a new paragraph to the text frame

# Shape Types
Shape.auto_shape_type  # Gets or sets the type of an auto shape

# Setting Slide Background
Slide.background.fill.solid()  # Sets the slide background fill to a solid color
